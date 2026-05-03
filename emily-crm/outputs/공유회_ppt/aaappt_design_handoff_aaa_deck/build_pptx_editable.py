"""
AAA 발표 덱 v7 — 편집 가능 PPTX 빌더
이미지 기반(build_pptx.py)과 달리 텍스트박스·도형으로 재구성해 PowerPoint에서 직접 편집 가능.
디자인 토큰은 deck-styles.css / slides-*.jsx 기준.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

HERE = Path(__file__).parent
OUT = HERE / "AAA_발표_덱_v7_editable.pptx"
ASSETS = HERE / "assets"

# ── Design tokens ──────────────────────────────────────────
INK = RGBColor(0x0A, 0x0A, 0x0A)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
YELLOW = RGBColor(0xE9, 0xED, 0x12)
YELLOW_INK = RGBColor(0x3A, 0x3B, 0x00)
MUTE = RGBColor(0x88, 0x88, 0x88)
MUTE2 = RGBColor(0xC9, 0xC9, 0xC9)
SURFACE = RGBColor(0xF5, 0xF5, 0xF5)
RULE = RGBColor(0xE6, 0xE6, 0xE6)
CODE_BG = RGBColor(0x11, 0x11, 0x11)
CODE_FG = RGBColor(0xEA, 0xEA, 0xEA)
KW = RGBColor(0xFF, 0x7E, 0x8A)
ST = RGBColor(0xFF, 0xD6, 0x7A)
CM = RGBColor(0x6A, 0x6A, 0x6A)

FONT_KO = "Pretendard Variable"
FONT_MONO = "JetBrains Mono"

# ── px → EMU ───────────────────────────────────────────────
def px(n): return Emu(int(n * 9525))

SLIDE_W, SLIDE_H = 1920, 1080

# ── Helpers ────────────────────────────────────────────────
def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None, radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, px(x), px(y), px(w), px(h))
    if radius and shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
        # Adjust corner radius (0.0 ~ 0.5)
        try:
            shp.adjustments[0] = min(0.5, radius / min(w, h))
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_w:
            shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp

def add_text(slide, x, y, w, h, text, *, size=24, bold=False, color=INK,
             align='left', font=FONT_KO, letter_spacing=None, line_height=None,
             anchor='top'):
    tb = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    anchor_map = {'top': MSO_ANCHOR.TOP, 'middle': MSO_ANCHOR.MIDDLE, 'bottom': MSO_ANCHOR.BOTTOM}
    tf.vertical_anchor = anchor_map.get(anchor, MSO_ANCHOR.TOP)

    lines = text.split('\n') if isinstance(text, str) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}
        p.alignment = align_map.get(align, PP_ALIGN.LEFT)
        if line_height:
            p.line_spacing = line_height
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        if letter_spacing is not None:
            _set_letter_spacing(r, letter_spacing)
    return tb

def add_runs(slide, x, y, w, h, runs, *, align='left', line_height=None, anchor='top'):
    """runs: list of dicts with keys text, size, bold, color, font, highlight, underline, letter_spacing"""
    tb = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    anchor_map = {'top': MSO_ANCHOR.TOP, 'middle': MSO_ANCHOR.MIDDLE, 'bottom': MSO_ANCHOR.BOTTOM}
    tf.vertical_anchor = anchor_map.get(anchor, MSO_ANCHOR.TOP)
    align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}

    paragraphs = [[]]
    for r in runs:
        t = r.get('text', '')
        if '\n' in t:
            parts = t.split('\n')
            for j, part in enumerate(parts):
                if j > 0:
                    paragraphs.append([])
                if part:
                    paragraphs[-1].append({**r, 'text': part})
        else:
            paragraphs[-1].append(r)

    for pi, p_runs in enumerate(paragraphs):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align_map.get(align, PP_ALIGN.LEFT)
        if line_height:
            p.line_spacing = line_height
        if not p_runs:
            # empty paragraph — blank line
            continue
        for ri, r_data in enumerate(p_runs):
            run = p.add_run()
            run.text = r_data.get('text', '')
            run.font.name = r_data.get('font', FONT_KO)
            run.font.size = Pt(r_data.get('size', 24))
            run.font.bold = r_data.get('bold', False)
            run.font.color.rgb = r_data.get('color', INK)
            if r_data.get('underline'):
                run.font.underline = True
            if r_data.get('highlight'):
                _set_highlight(run, r_data['highlight'])
            if r_data.get('letter_spacing') is not None:
                _set_letter_spacing(run, r_data['letter_spacing'])
    return tb

def _set_highlight(run, rgb):
    rPr = run._r.get_or_add_rPr()
    # Remove existing highlight
    for el in rPr.findall(qn('a:highlight')):
        rPr.remove(el)
    hl = etree.SubElement(rPr, qn('a:highlight'))
    solidFill = etree.SubElement(hl, qn('a:solidFill'))
    clr = etree.SubElement(solidFill, qn('a:srgbClr'))
    clr.set('val', '{:02X}{:02X}{:02X}'.format(rgb[0], rgb[1], rgb[2]))

def _set_letter_spacing(run, em):
    # em unit to PowerPoint spc (1/100 of a point)
    # rough: -0.02em on 100pt text = -2pt → spc = -200
    rPr = run._r.get_or_add_rPr()
    spc_val = int(em * 1000)
    rPr.set('spc', str(spc_val))

def add_top_bar(slide, left, right, *, dark=False, yellow=False):
    # 옐로우 사각 인디케이터 + 좌우 메타 텍스트
    meta_color = RGBColor(0x88, 0x88, 0x88)
    if dark:
        meta_color = RGBColor(0x66, 0x66, 0x66)
    if yellow:
        meta_color = RGBColor(0x6A, 0x68, 0x00)
    # indicator square
    ind_color = YELLOW if not yellow else INK
    add_rect(slide, 64, 48, 10, 10, fill=ind_color)
    # left label
    add_text(slide, 90, 40, 600, 30, left or "", size=14, bold=True, color=meta_color,
             letter_spacing=0.18, font=FONT_KO)
    # right label
    add_text(slide, SLIDE_W - 800 - 64, 40, 800, 30, right or "EMILY · SELFISH CLUB · AAA",
             size=14, bold=True, color=meta_color, letter_spacing=0.18, align='right')

def add_bot_bar(slide, left, right, *, dark=False, yellow=False):
    meta_color = RGBColor(0x88, 0x88, 0x88)
    if dark:
        meta_color = RGBColor(0x66, 0x66, 0x66)
    if yellow:
        meta_color = RGBColor(0x6A, 0x68, 0x00)
    if left:
        add_text(slide, 64, SLIDE_H - 60, 800, 30, left, size=14, bold=True, color=meta_color,
                 letter_spacing=0.18)
    if right:
        add_text(slide, SLIDE_W - 800 - 64, SLIDE_H - 60, 800, 30, right, size=14, bold=True,
                 color=meta_color, letter_spacing=0.18, align='right')

def add_kicker(slide, x, y, text, *, dark=False):
    color = RGBColor(0x77, 0x77, 0x77) if not dark else RGBColor(0x88, 0x88, 0x88)
    add_text(slide, x, y, 1400, 40, text, size=18, bold=True, color=color,
             letter_spacing=0.22)

def add_picture_safe(slide, img_path, x, y, w, h):
    if not img_path.exists():
        return None
    return slide.shapes.add_picture(str(img_path), px(x), px(y), px(w), px(h))

# ── Slide builders ────────────────────────────────────────
def new_slide(prs, bg=PAPER):
    blank = prs.slide_layouts[6]
    s = prs.slides.add_slide(blank)
    set_bg(s, bg)
    return s

def s01(prs):
    s = new_slide(prs, PAPER)
    add_top_bar(s, "INTRO", "EMILY · SELFISH CLUB · AAA")
    # 매일 아침,
    add_text(s, 140, 340, 1600, 180, "매일 아침,", size=110, bold=False, color=INK,
             letter_spacing=-0.04, line_height=1.12)
    # AI가 [알림톡]을 — 하이라이트
    add_runs(s, 140, 470, 1600, 180, [
        {'text': 'AI가 ', 'size': 110, 'bold': True, 'color': INK, 'letter_spacing': -0.04},
        {'text': '알림톡', 'size': 110, 'bold': True, 'color': INK, 'highlight': YELLOW, 'letter_spacing': -0.04},
        {'text': '을', 'size': 110, 'bold': True, 'color': INK, 'letter_spacing': -0.04},
    ], line_height=1.12)
    # 들고 옵니다.
    add_text(s, 140, 600, 1600, 180, "들고 옵니다.", size=110, bold=False, color=INK,
             letter_spacing=-0.04, line_height=1.12)
    # 메타
    add_text(s, 140, 870, 1600, 40, "에밀리 · 셀피쉬클럽 CRM PM   ·   AAA 공유회 · 2026.04.28",
             size=20, color=RGBColor(0xA8, 0xA8, 0xA8), letter_spacing=-0.01)
    return s

def s02(prs):
    s = new_slide(prs, PAPER)
    add_top_bar(s, "INTRO · ABOUT", "EMILY · SELFISH CLUB · AAA")
    # 에밀리.
    add_runs(s, 140, 220, 900, 220, [
        {'text': '에밀리', 'size': 140, 'bold': True, 'color': INK, 'letter_spacing': -0.04},
        {'text': '.', 'size': 140, 'bold': True, 'color': YELLOW, 'letter_spacing': -0.04},
    ], line_height=1.0)
    # 역할
    add_text(s, 140, 470, 900, 65, "셀피쉬클럽 CRM PM", size=32, bold=False, color=INK,
             letter_spacing=-0.02, line_height=1.35)
    add_text(s, 140, 530, 900, 65, "인하우스 풀스택 마케터", size=32, bold=False, color=INK,
             letter_spacing=-0.02, line_height=1.35)
    # 태그 5개
    tags = [('퍼포먼스', False), ('그로스', False), ('ops', False),
            ('효율충', True), ('데이터자동화', False)]
    tx = 140
    ty = 660
    for t, hl in tags:
        text = f"#{t}"
        # 대략 글자폭 계산
        w_px = 24 + len(text) * 22
        if hl:
            add_rect(s, tx, ty, w_px, 54, fill=YELLOW, radius=27)
            add_text(s, tx, ty + 14, w_px, 30, text, size=18, color=YELLOW_INK, align='center')
        else:
            add_rect(s, tx, ty, w_px, 54, fill=None, line=RULE, line_w=1.5, radius=27)
            add_text(s, tx, ty + 14, w_px, 30, text, size=18, color=RGBColor(0x66, 0x66, 0x66), align='center')
        tx += w_px + 12
    # 프로필 카드
    card_x, card_y, cw, ch = 1290, 220, 480, 580
    add_rect(s, card_x, card_y, cw, ch, fill=SURFACE, radius=24)
    add_picture_safe(s, ASSETS / 'emily.png', card_x, card_y, cw, ch)
    add_text(s, card_x + 24, card_y + ch - 50, 200, 30, "EMILY", size=14, bold=True,
             color=PAPER, letter_spacing=0.24)
    return s

def s03(prs):
    s = new_slide(prs, PAPER)
    add_top_bar(s, "INTRO · AAA", "EMILY · SELFISH CLUB · AAA")
    add_kicker(s, 140, 180, "지난 6주, AAA팀에 있었어요.")
    add_runs(s, 140, 240, 1640, 140, [
        {'text': '6주 동안 ', 'size': 78, 'bold': True, 'color': INK, 'letter_spacing': -0.04},
        {'text': '뭘', 'size': 78, 'bold': True, 'color': INK, 'highlight': YELLOW, 'letter_spacing': -0.04},
        {'text': ' 했냐면.', 'size': 78, 'bold': True, 'color': INK, 'letter_spacing': -0.04},
    ], line_height=1.06)
    # 3 카드
    cards = [
        ('WHY',  '반복에 쓰는 시간을\n기획에 쓰고 싶어서.'),
        ('WHAT', 'Claude Code로\nCRM 자동화를 직접.'),
        ('NOW',  '아침 10시,\n슬랙 버튼 하나로 끝.'),
    ]
    cx, cy, cw, ch, gap = 140, 500, 540, 360, 20
    for i, (k, t) in enumerate(cards):
        x = cx + i * (cw + gap)
        add_rect(s, x, cy, cw, ch, fill=SURFACE, radius=16)
        add_text(s, x + 44, cy + 48, cw - 88, 30, k, size=14, bold=True,
                 color=RGBColor(0x88, 0x88, 0x88), letter_spacing=0.18)
        add_text(s, x + 44, cy + 100, cw - 88, ch - 120, t, size=30, bold=False, color=INK,
                 line_height=1.3, letter_spacing=-0.03)
    return s

def s04(prs):
    s = new_slide(prs, PAPER)
    add_top_bar(s, "CONTENTS", "EMILY · SELFISH CLUB · AAA")
    add_runs(s, 140, 180, 1640, 120, [
        {'text': '목차', 'size': 90, 'bold': True, 'color': INK, 'letter_spacing': -0.04},
        {'text': '.', 'size': 90, 'bold': True, 'color': YELLOW, 'letter_spacing': -0.04},
    ])
    rows = [
        ('01', '지난 여정',  '예전에 제가 일했던 방식'),
        ('02', 'Step 1',   '자동화 구조 설계'),
        ('03', 'Step 2',   'AI한테 내 맥락 심기 · 하네스'),
        ('04', '마무리',    '팀과 일했던 이야기'),
    ]
    col_w = 820
    row_h = 140
    top = 400
    for i, (n, k, t) in enumerate(rows):
        r, c = i // 2, i % 2
        x = 140 + c * (col_w + 20)
        y = top + r * (row_h + 24)
        # border top
        add_rect(s, x, y, col_w, 2, fill=RULE)
        add_text(s, x, y + 30, 160, 80, n, size=54, color=RGBColor(0xBB, 0xBB, 0xBB),
                 letter_spacing=-0.04)
        add_text(s, x + 180, y + 30, col_w - 200, 30, k.upper(), size=14, bold=True,
                 color=RGBColor(0x88, 0x88, 0x88), letter_spacing=0.22)
        add_text(s, x + 180, y + 60, col_w - 200, 60, t, size=26, bold=True, color=INK,
                 letter_spacing=-0.02)
    return s

def s05(prs):
    s = new_slide(prs, PAPER)
    add_top_bar(s, "TL;DR", "EMILY · SELFISH CLUB · AAA")
    add_kicker(s, 140, 180, "오늘 제 세션이 끝나면,")
    add_runs(s, 140, 240, 1640, 130, [
        {'text': '가져가실 수 있는 3가지', 'size': 78, 'bold': True, 'color': INK,
         'highlight': YELLOW, 'letter_spacing': -0.04},
    ])
    rows = [
        ('01', 'AI랑 같이 기획하는 법'),
        ('02', '자동화 구조 설계'),
        ('03', '내 규칙을 AI가 읽게 만드는 법'),
    ]
    top = 440
    rh = 150
    for i, (n, t) in enumerate(rows):
        y = top + i * rh
        add_rect(s, 140, y, 1640, 2, fill=RULE)
        add_text(s, 140, y + 30, 160, 90, n, size=60, color=RGBColor(0xBB, 0xBB, 0xBB),
                 letter_spacing=-0.04)
        add_text(s, 320, y + 30, 1460, 90, t, size=48, bold=True, color=INK,
                 letter_spacing=-0.02, anchor='top')
    return s

def channel_card(slide, x, y, n, label, img, brand_color, brand_text):
    card_w = 260
    card_h = 520
    # 프레임 (검정 베젤)
    add_rect(slide, x, y, card_w, card_h, fill=RGBColor(0x1A, 0x1A, 0x1A), radius=36)
    # 노치
    add_rect(slide, x + card_w // 2 - 40, y + 10, 80, 20, fill=RGBColor(0, 0, 0), radius=999)
    # 스크린 (흰 배경)
    sw, sh = card_w - 16, card_h - 16
    add_rect(slide, x + 8, y + 8, sw, sh, fill=PAPER, radius=28)
    # 스크린 내부 이미지 (노치 아래부터)
    if img and img.exists():
        add_picture_safe(slide, img, x + 12, y + 60, sw - 8, sh - 70)
    # 번호 + 배지 + 라벨 (상단)
    add_text(slide, x, y - 40, 40, 24, n, size=12, color=RGBColor(0xBB, 0xBB, 0xBB))
    add_rect(slide, x + 36, y - 44, 22, 22, fill=brand_color, radius=5)
    add_text(slide, x + 36, y - 42, 22, 20, brand_text, size=10, bold=True,
             color=RGBColor(0x3C, 0x1E, 0x1E) if brand_color == YELLOW else PAPER,
             align='center')
    add_text(slide, x + 66, y - 40, card_w - 66, 24, label.upper(), size=12, bold=True,
             color=INK, letter_spacing=0.14)

def s06(prs):
    s = new_slide(prs, PAPER)
    add_top_bar(s, "§ 1 · 지난 여정", "EMILY · SELFISH CLUB · AAA")
    add_kicker(s, 140, 160, "근데 그 전에.")
    add_runs(s, 140, 210, 1640, 160, [
        {'text': '여러분은 어떤 ', 'size': 44, 'color': INK, 'letter_spacing': -0.03},
        {'text': '채널·메시지', 'size': 44, 'bold': True, 'color': INK, 'highlight': YELLOW,
         'letter_spacing': -0.03},
        {'text': '를 통해\n이 공유회에 오게 되셨나요?', 'size': 44, 'color': INK,
         'letter_spacing': -0.03},
    ], line_height=1.25)
    # 5개 카드
    labels = [
        ('01', '알림톡',  'alimtalk.png',   YELLOW, 'K'),
        ('02', '카플친',  'kakaoplus.png',  YELLOW, 'K'),
        ('03', '인스타',  'insta_main.png', RGBColor(0xDC, 0x27, 0x43), '◉'),
        ('04', '오카방',  'oka.png',        RGBColor(0xB2, 0xC7, 0xD9), '#'),
        ('05', '이메일',  'Email.png',      PAPER, 'M'),
    ]
    x0 = 160
    gap = 24
    cw = 260
    y0 = 500
    for i, (n, label, img, bc, bt) in enumerate(labels):
        x = x0 + i * (cw + gap)
        channel_card(s, x, y0, n, label, ASSETS / img, bc, bt)
    return s

def s07(prs):
    s = new_slide(prs, PAPER)
    add_top_bar(s, "§ 1 · 지난 여정", "EMILY · SELFISH CLUB · AAA")
    add_kicker(s, 140, 160, "공유회 1건이 열리면,")
    add_runs(s, 140, 210, 1640, 160, [
        {'text': '채널 ', 'size': 78, 'bold': True, 'color': INK, 'letter_spacing': -0.04},
        {'text': '네 개', 'size': 78, 'bold': True, 'color': INK, 'underline': True,
         'letter_spacing': -0.04},
        {'text': '가 동시에 나가요.', 'size': 78, 'bold': True, 'color': INK,
         'letter_spacing': -0.04},
    ])
    # 4 카드
    cards = [
        ('01', '알림톡', '카카오 · 템플릿 9종'),
        ('02', '카플친', '플러스친구 소식'),
        ('03', '채널톡', '앱 내 팝업'),
        ('04', '오카방', '오픈채팅 공지'),
    ]
    x0 = 140
    cw = 400
    gap = 20
    y0 = 460
    for i, (n, t, d) in enumerate(cards):
        x = x0 + i * (cw + gap)
        add_rect(s, x, y0, cw, 280, fill=SURFACE, radius=16)
        add_text(s, x + 36, y0 + 40, cw - 72, 30, n, size=14, color=RGBColor(0xBB, 0xBB, 0xBB))
        add_text(s, x + 36, y0 + 80, cw - 72, 80, t, size=40, bold=True, color=INK,
                 letter_spacing=-0.03)
        add_text(s, x + 36, y0 + 220, cw - 72, 40, d, size=16, color=RGBColor(0x77, 0x77, 0x77))
    # 하단 설명
    add_runs(s, 140, 820, 1640, 60, [
        {'text': '오늘은 이 중에서 ', 'size': 22, 'color': RGBColor(0x44, 0x44, 0x44)},
        {'text': '알림톡', 'size': 22, 'bold': True, 'color': INK, 'highlight': YELLOW},
        {'text': ' 얘기만 할 거예요. 제일 빡센 녀석이거든요.', 'size': 22,
         'color': RGBColor(0x44, 0x44, 0x44)},
    ])
    return s

def s08(prs):
    s = new_slide(prs, PAPER)
    add_top_bar(s, "§ 1 · 지난 여정", "EMILY · SELFISH CLUB · AAA")
    add_kicker(s, 140, 160, "알림톡, 얼마나 빡세냐면.")
    add_text(s, 140, 210, 1640, 100, "공유회 1건당 이렇게 나가요.", size=48, bold=False,
             color=INK, letter_spacing=-0.035)
    # 3 big cards
    cards = [
        ('10', '일',          'D-10 부터 D-day 까지', False),
        ('9',  '번 발송',       '승인 템플릿 9종',      False),
        ('N',  '번 바뀌는 대상자', '매 시점마다 다름',   True),  # yellow
    ]
    x0 = 140
    cw = 533
    gap = 20
    y0 = 380
    for i, (num, unit, t, is_y) in enumerate(cards):
        x = x0 + i * (cw + gap)
        if is_y:
            add_rect(s, x, y0, cw, 400, fill=INK, radius=16)
            num_c = YELLOW
            unit_c = RGBColor(0xAA, 0xAA, 0xAA)
            desc_c = RGBColor(0xDD, 0xDD, 0xDD)
        else:
            add_rect(s, x, y0, cw, 400, fill=SURFACE, radius=16)
            num_c = INK
            unit_c = RGBColor(0x66, 0x66, 0x66)
            desc_c = RGBColor(0x33, 0x33, 0x33)
        add_text(s, x + 48, y0 + 56, 200, 180, num, size=140, bold=True, color=num_c,
                 letter_spacing=-0.05)
        add_text(s, x + 260, y0 + 120, cw - 260, 60, unit, size=32, bold=True, color=unit_c)
        add_text(s, x + 48, y0 + 280, cw - 96, 60, t, size=20, color=desc_c)
    # 하단 모노
    add_runs(s, 140, 830, 1640, 40, [
        {'text': '· D-5 ', 'size': 14, 'color': RGBColor(0x88, 0x88, 0x88), 'font': FONT_MONO},
        {'text': '오픈알림 안 받은 사람만', 'size': 14, 'color': RGBColor(0x33, 0x33, 0x33),
         'font': FONT_MONO},
        {'text': '    · D-1 ', 'size': 14, 'color': RGBColor(0x88, 0x88, 0x88),
         'font': FONT_MONO},
        {'text': '신청자만', 'size': 14, 'color': RGBColor(0x33, 0x33, 0x33),
         'font': FONT_MONO},
        {'text': '    · D-day ', 'size': 14, 'color': RGBColor(0x88, 0x88, 0x88),
         'font': FONT_MONO},
        {'text': '결제자만', 'size': 14, 'color': RGBColor(0x33, 0x33, 0x33),
         'font': FONT_MONO},
    ])
    return s

def s09(prs):
    s = new_slide(prs, PAPER)
    add_top_bar(s, "§ 1 · 지난 여정", "EMILY · SELFISH CLUB · AAA")
    # 왼쪽
    add_kicker(s, 140, 280, "근데 이거, 생각보다 손이 많이 가요.")
    add_runs(s, 140, 330, 680, 280, [
        {'text': '매번\n', 'size': 52, 'bold': True, 'color': INK, 'letter_spacing': -0.035},
        {'text': '수동으로 DB', 'size': 52, 'bold': True, 'color': INK, 'underline': True,
         'letter_spacing': -0.035},
        {'text': '를\n발랐어요.', 'size': 52, 'bold': True, 'color': INK,
         'letter_spacing': -0.035},
    ], line_height=1.2)
    add_text(s, 140, 650, 680, 120,
             "신청자 추출 → 미신청자 분리 → 다음 타겟 세팅.\n조건 하나 바뀌면, 또 또 또.",
             size=20, color=RGBColor(0x44, 0x44, 0x44), line_height=1.5)
    # 코드블록
    code_x, code_y, code_w, code_h = 900, 220, 880, 700
    add_rect(s, code_x, code_y, code_w, code_h, fill=CODE_BG, radius=12)
    code_lines = [
        ('-- 공유회 오픈 알림 대상자', CM),
        ('SELECT', KW, ' u_phone, u_name', CODE_FG),
        ('FROM', KW, ' membership', CODE_FG),
        ('WHERE', KW, ' marketing_agree = ', CODE_FG, 'true', ST, ';', CODE_FG),
        ('', CODE_FG),
        ('-- 리마인드 (미신청자만)', CM),
        ('SELECT', KW, ' u_phone ', CODE_FG, 'FROM', KW, ' membership m', CODE_FG),
        ('WHERE NOT EXISTS', KW, ' (', CODE_FG),
        ('  SELECT', KW, ' 1 ', CODE_FG, 'FROM', KW, ' event e', CODE_FG),
        ('  WHERE', KW, ' e.u_phone = m.u_phone);', CODE_FG),
        ('', CODE_FG),
        ('-- D-1 확정자만 골라내기', CM),
        ('SELECT', KW, ' u_phone ', CODE_FG, 'FROM', KW, ' event', CODE_FG),
        ('WHERE', KW, ' status = ', CODE_FG, "'confirmed'", ST, ';', CODE_FG),
    ]
    # 간단히 모노스페이스 전체 텍스트박스 1개로
    line_y = code_y + 28
    for parts in code_lines:
        runs = []
        for i in range(0, len(parts), 2):
            runs.append({'text': parts[i], 'size': 14, 'color': parts[i + 1],
                         'font': FONT_MONO})
        if not any(r['text'] for r in runs):
            runs = [{'text': ' ', 'size': 14, 'color': CODE_FG, 'font': FONT_MONO}]
        add_runs(s, code_x + 32, line_y, code_w - 64, 28, runs)
        line_y += 30
    return s

def s10(prs):
    s = new_slide(prs, PAPER)
    add_top_bar(s, "§ 1 · 지난 여정", "EMILY · SELFISH CLUB · AAA")
    add_kicker(s, 140, 160, "그래서 AAA 6주간,")
    add_runs(s, 140, 210, 1640, 140, [
        {'text': '자동화', 'size': 90, 'bold': True, 'color': INK, 'highlight': YELLOW,
         'letter_spacing': -0.04},
        {'text': '를 시작했어요.', 'size': 90, 'bold': True, 'color': INK,
         'letter_spacing': -0.04},
    ])
    # BEFORE (SURFACE)
    bx, by, bw, bh = 140, 440, 760, 400
    add_rect(s, bx, by, bw, bh, fill=SURFACE, radius=16)
    add_text(s, bx + 44, by + 60, bw - 88, 30, 'BEFORE', size=14, bold=True,
             color=RGBColor(0x88, 0x88, 0x88), letter_spacing=0.18)
    add_text(s, bx + 44, by + 120, 200, 180, '3', size=160, bold=True,
             color=RGBColor(0xBB, 0xBB, 0xBB), letter_spacing=-0.05)
    add_text(s, bx + 280, by + 180, 150, 60, 'h', size=40, bold=True,
             color=RGBColor(0x99, 0x99, 0x99))
    add_text(s, bx + 44, by + 320, bw - 88, 50, '하루 3시간, 수동 SQL + 발송',
             size=22, color=RGBColor(0x33, 0x33, 0x33))
    # ARROW
    add_text(s, 920, 600, 80, 100, '→', size=64, color=RGBColor(0xBB, 0xBB, 0xBB),
             align='center')
    # AFTER (INK)
    ax, ay, aw, ah = 1020, 440, 760, 400
    add_rect(s, ax, ay, aw, ah, fill=INK, radius=16)
    add_text(s, ax + 44, ay + 60, aw - 88, 30, 'AFTER', size=14, bold=True,
             color=YELLOW, letter_spacing=0.18)
    add_text(s, ax + 44, ay + 120, 200, 180, '5', size=160, bold=True,
             color=YELLOW, letter_spacing=-0.05)
    add_text(s, ax + 280, ay + 180, 200, 60, 'min', size=40, bold=True,
             color=RGBColor(0xAA, 0xAA, 0xAA))
    add_text(s, ax + 44, ay + 320, aw - 88, 50, '슬랙 승인 한 번 누르면 끝.',
             size=22, color=RGBColor(0xDD, 0xDD, 0xDD))
    # 하단
    add_runs(s, 140, 880, 1640, 60, [
        {'text': '오늘 아침 제 슬랙엔 이게 와요. ', 'size': 20,
         'color': RGBColor(0x66, 0x66, 0x66)},
        {'text': '→ 다음 장', 'size': 20, 'bold': True, 'color': INK},
    ])
    return s

def slack_modal(slide, x, y):
    w, h = 620, 500
    add_rect(slide, x, y, w, h, fill=PAPER, line=RGBColor(0xE4, 0xE4, 0xE4), line_w=1, radius=14)
    # 헤더
    add_rect(slide, x, y, w, 60, fill=PAPER, radius=14)
    add_text(slide, x + 22, y + 18, 500, 30, '# 오늘의 알림톡 · 승인 요청', size=14, bold=True)
    add_text(slide, x + w - 40, y + 18, 20, 30, '✕', size=14,
             color=RGBColor(0x99, 0x99, 0x99))
    # 봇 정보 바
    add_rect(slide, x, y + 60, w, 70, fill=RGBColor(0xF9, 0xF9, 0xF9))
    add_rect(slide, x + 22, y + 77, 36, 36, fill=YELLOW, radius=8)
    add_text(slide, x + 22, y + 83, 36, 30, '🌟', size=16, align='center')
    add_text(slide, x + 72, y + 75, 300, 30, '알림톡봇  앱', size=12, bold=True)
    add_text(slide, x + 72, y + 100, 300, 30, '오전 11:05', size=10,
             color=RGBColor(0x88, 0x88, 0x88))
    # 본문 영역
    add_text(slide, x + 22, y + 150, w - 44, 40,
             '🔗 [#6 입장링크] 비개발자를 위한 클로드코드앱 & 깃헙 실전 정복기',
             size=12, bold=True)
    add_text(slide, x + 22, y + 180, w - 44, 30, '— 수정됨 ✏️', size=11,
             color=RGBColor(0x88, 0x88, 0x88))
    add_text(slide, x + 22, y + 210, w - 44, 30,
             '대상: 1명   ⏰ 발송시간: 11:11   🔗 URL: sepia-quartz-81f.notion.site/3305c...',
             size=10, color=RGBColor(0x44, 0x44, 0x44))
    # preview box
    add_rect(slide, x + 22, y + 250, w - 44, 120, fill=RGBColor(0xF5, 0xF5, 0xF5), radius=8)
    add_text(slide, x + 38, y + 262, w - 76, 30, '📋 수정된 미리보기:', size=11, bold=True)
    add_runs(slide, x + 38, y + 290, w - 76, 60, [
        {'text': '{이름}님께서 신청하신 ', 'size': 10, 'color': RGBColor(0x33, 0x33, 0x33)},
        {'text': '이기적공유회', 'size': 10, 'bold': True, 'color': RGBColor(0x33, 0x33, 0x33)},
        {'text': '가 ', 'size': 10, 'color': RGBColor(0x33, 0x33, 0x33)},
        {'text': '20:00', 'size': 10, 'bold': True, 'color': INK, 'highlight': YELLOW},
        {'text': '에 시작됩니다 😊 늦지 않게 입장해주세요!', 'size': 10,
         'color': RGBColor(0x33, 0x33, 0x33)},
    ], line_height=1.5)
    # 경고
    add_text(slide, x + 22, y + 385, w - 44, 30,
             '⚠️ 승인 시 예약 발송이 확정됩니다. 카피·링크·시간을 반드시 확인 후 승인해주세요.',
             size=9, color=RGBColor(0x99, 0x99, 0x99))
    # 버튼들
    add_rect(slide, x + 22, y + 420, 140, 40, fill=RGBColor(0x2E, 0xB6, 0x7D), radius=6)
    add_text(slide, x + 22, y + 430, 140, 30, '✅ 이대로 발송', size=11, bold=True,
             color=PAPER, align='center')
    add_rect(slide, x + 170, y + 420, 80, 40, fill=PAPER, line=RGBColor(0xCC, 0xCC, 0xCC),
             line_w=1, radius=6)
    add_text(slide, x + 170, y + 430, 80, 30, '✏️ 수정', size=11, bold=True, align='center')
    add_rect(slide, x + 258, y + 420, 80, 40, fill=PAPER, line=RGBColor(0xCC, 0xCC, 0xCC),
             line_w=1, radius=6)
    add_text(slide, x + 258, y + 430, 80, 30, '❌ 취소', size=11, bold=True,
             color=RGBColor(0xC4, 0x34, 0x3B), align='center')

def s11(prs):
    s = new_slide(prs, PAPER)
    add_top_bar(s, "§ 1 · 지난 여정", "EMILY · SELFISH CLUB · AAA")
    add_kicker(s, 140, 220, "지금 제 하루는 이렇게 시작해요. 슬랙 하나로.")
    add_runs(s, 140, 270, 800, 220, [
        {'text': '오늘 아침\n슬랙 한 번에 ', 'size': 54, 'bold': True, 'color': INK,
         'letter_spacing': -0.035},
        {'text': '끝', 'size': 54, 'bold': True, 'color': INK, 'highlight': YELLOW,
         'letter_spacing': -0.035},
        {'text': '.', 'size': 54, 'bold': True, 'color': INK, 'letter_spacing': -0.035},
    ], line_height=1.2)
    # 3 items
    items = [
        ('①', 'AM 10:00 · 모달이 와요', '오늘 나갈 메시지 통합 리스트'),
        ('②', '수정·승인, 한 번에 끝', '시간·링크·카피 다 모달 안에서'),
        ('③', '예약 발송, 알아서 나가요', 'Solapi로 자동'),
    ]
    iy = 560
    for n, t, d in items:
        add_text(s, 140, iy, 60, 40, n, size=22, bold=True, color=RGBColor(0xBB, 0xBB, 0xBB))
        add_text(s, 210, iy, 700, 40, t, size=20, bold=True, color=INK)
        add_text(s, 210, iy + 36, 700, 30, d, size=15, color=RGBColor(0x66, 0x66, 0x66))
        iy += 80
    # 슬랙 모달
    slack_modal(s, 1100, 250)
    return s

def s12(prs):
    s = new_slide(prs, INK)
    add_top_bar(s, "§ 2 · STEP 1", "EMILY · SELFISH CLUB · AAA", dark=True)
    add_text(s, 140, 180, 1640, 40, "여기서부터 본편.", size=18, bold=True,
             color=RGBColor(0x99, 0x99, 0x99), letter_spacing=0.22)
    add_runs(s, 140, 250, 1640, 500, [
        {'text': 'Step 1.\n', 'size': 150, 'bold': True, 'color': PAPER,
         'letter_spacing': -0.04},
        {'text': '자동화 구조', 'size': 150, 'bold': True, 'color': PAPER, 'underline': True,
         'letter_spacing': -0.04},
        {'text': '\n설계.', 'size': 150, 'bold': True, 'color': PAPER,
         'letter_spacing': -0.04},
    ], line_height=1.02)
    add_text(s, 140, 870, 1640, 40,
             "노트북을 덮어도, 퇴근한 뒤에도 — 메시지가 나가게.", size=24,
             color=RGBColor(0xBB, 0xBB, 0xBB))
    return s

def s13(prs):
    s = new_slide(prs, INK)
    add_top_bar(s, "§ 2 · STEP 1", "EMILY · SELFISH CLUB · AAA", dark=True)
    add_text(s, 140, 180, 1640, 40, "그 전에,", size=18, bold=True,
             color=RGBColor(0x99, 0x99, 0x99), letter_spacing=0.22)
    add_runs(s, 140, 230, 1640, 140, [
        {'text': '도구 ', 'size': 78, 'bold': True, 'color': PAPER,
         'letter_spacing': -0.04},
        {'text': '셋', 'size': 78, 'bold': True, 'color': INK, 'highlight': YELLOW,
         'letter_spacing': -0.04},
        {'text': '.', 'size': 78, 'bold': True, 'color': PAPER, 'letter_spacing': -0.04},
    ])
    add_text(s, 140, 390, 1640, 40, "셀피쉬가 쓰는 스택이에요.", size=22,
             color=RGBColor(0xD0, 0xD0, 0xD0))
    cards = [
        ('01', 'Slack',    '커뮤니케이션 · 승인', '메신저이자 승인 게이트'),
        ('02', 'Supabase', '데이터베이스',        '멤버·신청·일정 한 곳에'),
        ('03', 'n8n',      '자동화 워크플로우',    '서버에서 24/7'),
    ]
    x0 = 140
    cw = 540
    gap = 20
    y0 = 500
    for i, (n, t, d1, d2) in enumerate(cards):
        x = x0 + i * (cw + gap)
        add_rect(s, x, y0, cw, 380, fill=RGBColor(0x1A, 0x1A, 0x1A),
                 line=RGBColor(0x33, 0x33, 0x33), line_w=1, radius=16)
        add_text(s, x + 48, y0 + 56, cw - 96, 30, n, size=14, bold=True,
                 color=RGBColor(0x77, 0x77, 0x77))
        add_text(s, x + 48, y0 + 96, cw - 96, 80, t, size=48, bold=True, color=PAPER,
                 letter_spacing=-0.03)
        add_text(s, x + 48, y0 + 220, cw - 96, 30, d1.upper(), size=14, bold=True,
                 color=RGBColor(0x88, 0x88, 0x88), letter_spacing=0.18)
        add_text(s, x + 48, y0 + 260, cw - 96, 60, d2, size=18,
                 color=RGBColor(0xDD, 0xDD, 0xDD))
    return s

def s14(prs):
    s = new_slide(prs, PAPER)
    add_top_bar(s, "§ 2 · STEP 1", "EMILY · SELFISH CLUB · AAA")
    add_bot_bar(s, "T — 1주 · 첫 삽질", "")
    add_kicker(s, 140, 160, "T — 1주 · 첫 삽질")
    add_runs(s, 140, 210, 1640, 140, [
        {'text': '노트북 덮으면 ', 'size': 54, 'bold': True, 'color': INK,
         'letter_spacing': -0.035},
        {'text': '멈추는', 'size': 54, 'bold': True, 'color': INK, 'highlight': YELLOW,
         'letter_spacing': -0.035},
        {'text': ' 자동화?', 'size': 54, 'bold': True, 'color': INK,
         'letter_spacing': -0.035},
    ])
    # chat bubbles
    chats = [
        ('나', 'L', '미신청자 뽑아서, 저녁 7시에 이거 보내줘.'),
        ('CLAUDE', 'R', '네, 7시에 보낼게요.'),
        ('나', 'L', '…근데 저 6시에 퇴근하거든요. 노트북 끄고 가도 돼요?'),
        ('CLAUDE', 'R', '아… 노트북이 꺼지면, 저는 못 보내요.'),
    ]
    cy = 400
    for who, side, t in chats:
        bubble_w = 800
        bubble_x = 140 if side == 'L' else SLIDE_W - 140 - bubble_w
        add_text(s, bubble_x, cy, bubble_w, 24, who, size=12, bold=True,
                 color=RGBColor(0x99, 0x99, 0x99), letter_spacing=0.16,
                 align='left' if side == 'L' else 'right')
        bg = RGBColor(0xF1, 0xF1, 0xF1) if side == 'L' else INK
        fg = INK if side == 'L' else PAPER
        add_rect(s, bubble_x, cy + 30, bubble_w, 80, fill=bg, radius=18)
        add_text(s, bubble_x + 28, cy + 52, bubble_w - 56, 40, t, size=20, color=fg)
        cy += 120
    add_runs(s, 140, 900, 1640, 60, [
        {'text': '→ 내가 없어도 돌아갈 ', 'size': 22, 'color': RGBColor(0x33, 0x33, 0x33)},
        {'text': '무언가', 'size': 22, 'bold': True, 'color': INK, 'highlight': YELLOW},
        {'text': '가 필요했어요.', 'size': 22, 'color': RGBColor(0x33, 0x33, 0x33)},
    ])
    return s

def s15(prs):
    s = new_slide(prs, PAPER)
    add_top_bar(s, "§ 2 · STEP 1", "EMILY · SELFISH CLUB · AAA")
    add_kicker(s, 140, 160, "그래서,")
    add_runs(s, 140, 210, 1640, 140, [
        {'text': '만드는 건 ', 'size': 78, 'bold': True, 'color': INK,
         'letter_spacing': -0.04},
        {'text': 'Claude', 'size': 78, 'bold': True, 'color': INK, 'highlight': YELLOW,
         'letter_spacing': -0.04},
        {'text': ', 보내는 건 ', 'size': 78, 'bold': True, 'color': INK,
         'letter_spacing': -0.04},
        {'text': 'n8n', 'size': 78, 'bold': True, 'color': INK, 'highlight': YELLOW,
         'letter_spacing': -0.04},
        {'text': '.', 'size': 78, 'bold': True, 'color': INK,
         'letter_spacing': -0.04},
    ])
    # 2 카드
    cards = [
        ('만드는 역할', 'Claude Code',
         '알림톡 내용·조건·타이밍을 설계.\n내가 노트북 켰을 때만 돌아요.', False),
        ('보내는 역할', 'n8n',
         '설계된 알림톡을 실제로 발송.\n서버에서 24시간 계속 돌아요.', True),
    ]
    cy, ch = 500, 380
    cw = 810
    gap = 20
    for i, (tag, t, d, is_k) in enumerate(cards):
        x = 140 + i * (cw + gap)
        bg = INK if is_k else SURFACE
        tag_bg = YELLOW if is_k else INK
        tag_fg = INK if is_k else YELLOW
        body_c = RGBColor(0xDD, 0xDD, 0xDD) if is_k else RGBColor(0x33, 0x33, 0x33)
        t_c = PAPER if is_k else INK
        add_rect(s, x, cy, cw, ch, fill=bg, radius=16)
        add_rect(s, x + 48, cy + 48, 150, 34, fill=tag_bg, radius=4)
        add_text(s, x + 48, cy + 54, 150, 24, tag, size=12, bold=True, color=tag_fg,
                 letter_spacing=0.18, align='center')
        add_text(s, x + 48, cy + 110, cw - 96, 80, t, size=48, bold=True, color=t_c,
                 letter_spacing=-0.03)
        add_text(s, x + 48, cy + 220, cw - 96, 120, d, size=20, color=body_c,
                 line_height=1.5)
    return s

def s16(prs):
    s = new_slide(prs, PAPER)
    add_top_bar(s, "§ 2 · STEP 1", "EMILY · SELFISH CLUB · AAA")
    add_kicker(s, 140, 160, "근데 또 불안했어요.")
    add_runs(s, 140, 210, 1640, 140, [
        {'text': '완전 자동은 ', 'size': 50, 'bold': True, 'color': INK,
         'letter_spacing': -0.035},
        {'text': '무서운', 'size': 50, 'bold': True, 'color': INK, 'highlight': YELLOW,
         'letter_spacing': -0.035},
        {'text': ' 이유가 있죠.', 'size': 50, 'bold': True, 'color': INK,
         'letter_spacing': -0.035},
    ])
    cards = [
        ('오타 한 번', '= 6,800명이 동시에 본다'),
        ('링크 깨짐',  '= 유입 0'),
        ('타겟 실수', '= 신청 안 한 사람한테 "내일 봐요"'),
    ]
    x0 = 140
    cw = 540
    gap = 20
    y0 = 440
    for i, (t, d) in enumerate(cards):
        x = x0 + i * (cw + gap)
        add_rect(s, x, y0, cw, 220, fill=SURFACE, radius=16)
        add_text(s, x + 36, y0 + 40, cw - 72, 50, t, size=28, bold=True, color=INK)
        add_text(s, x + 36, y0 + 110, cw - 72, 80, d, size=18,
                 color=RGBColor(0x55, 0x55, 0x55))
    # 하단 노란 박스
    add_rect(s, 140, 760, 1200, 100, fill=YELLOW, radius=14)
    add_runs(s, 170, 790, 1170, 60, [
        {'text': '사람이 한 번 보는 ', 'size': 26, 'bold': True, 'color': INK},
        {'text': '관문', 'size': 26, 'bold': True, 'color': INK, 'underline': True},
        {'text': '을 넣자. → Slack 승인.', 'size': 26, 'bold': True, 'color': INK},
    ])
    return s

def s17(prs):
    s = new_slide(prs, INK)
    add_top_bar(s, "§ 3 · STEP 2", "EMILY · SELFISH CLUB · AAA", dark=True)
    add_text(s, 140, 180, 1640, 40, "자동화만으론 안 되더라.", size=18, bold=True,
             color=RGBColor(0x99, 0x99, 0x99), letter_spacing=0.22)
    add_runs(s, 140, 250, 1640, 500, [
        {'text': 'Step 2.\n', 'size': 140, 'bold': True, 'color': PAPER,
         'letter_spacing': -0.04},
        {'text': 'AI한테 ', 'size': 140, 'bold': True, 'color': PAPER,
         'letter_spacing': -0.04},
        {'text': '내 맥락', 'size': 140, 'bold': True, 'color': INK, 'highlight': YELLOW,
         'letter_spacing': -0.04},
        {'text': '을\n심는 법.', 'size': 140, 'bold': True, 'color': PAPER,
         'letter_spacing': -0.04},
    ], line_height=1.02)
    add_runs(s, 140, 900, 1640, 50, [
        {'text': "'이게 ", 'size': 22, 'color': RGBColor(0xBB, 0xBB, 0xBB)},
        {'text': '하네스', 'size': 22, 'bold': True, 'color': YELLOW},
        {'text': "'구나'를 알기까지 좀 걸렸어요.", 'size': 22,
         'color': RGBColor(0xBB, 0xBB, 0xBB)},
    ])
    return s

def s18(prs):
    s = new_slide(prs, PAPER)
    add_top_bar(s, "§ 3 · STEP 2", "EMILY · SELFISH CLUB · AAA")
    # 왼쪽
    add_kicker(s, 140, 280, "규칙이 너무 많아서,")
    add_runs(s, 140, 330, 700, 220, [
        {'text': 'CLAUDE.md 353줄', 'size': 48, 'bold': True, 'color': INK,
         'underline': True, 'letter_spacing': -0.035},
        {'text': '에\n다 박았어요.', 'size': 48, 'bold': True, 'color': INK,
         'letter_spacing': -0.035},
    ], line_height=1.2)
    add_text(s, 140, 620, 700, 200,
             '근데 AI가 자꾸 틀리고 까먹음.\n맥북 2개 쓰는데, 다른 맥북 켜면 일주일 전 기억이 마지막.',
             size=20, color=RGBColor(0x44, 0x44, 0x44), line_height=1.5)
    # 코드블록 (오른쪽)
    cx, cy, cw, ch = 920, 240, 860, 680
    add_rect(s, cx, cy, cw, ch, fill=CODE_BG, radius=12)
    code_lines = [
        ('# 셀피쉬 CRM 운영 규칙', CM),
        ('# v2.1 — 2026.03.14', CM),
        ('', None),
        ('## 톤앤매너', YELLOW),
        ('- 존댓말 · 이모지 1개 이하', CODE_FG),
        ("- '~님' 호칭 통일", CODE_FG),
        ('- 느낌표 남발 금지', CODE_FG),
        ('', None),
        ('## 알림톡 타이밍', YELLOW),
        ('- D-5 오픈알림 / D-3 리마인드', CODE_FG),
        ('- D-1 확정자 / D-0 입장링크', CODE_FG),
        ('- D+1 후기 유도', CODE_FG),
        ('', None),
        ('## 타겟팅 예외', YELLOW),
        ('- 수신거부 자동 제외', CODE_FG),
        ('- 최근 30일 이탈자 제외', CODE_FG),
        ('- VIP는 별도 템플릿', CODE_FG),
        ('', None),
        ('## 링크 규칙', YELLOW),
        ('- 공유회 상세: event.selfishclub.co/:id', CODE_FG),
        ('- UTM: utm_source=kakao&utm_medium=alimtalk', CODE_FG),
        ('- 단축: s.slfc.to/ 항상 사용', CODE_FG),
        ('', None),
        ('## 변수 네이밍', YELLOW),
        ('- {{member_name}} / {{event_title}}', CODE_FG),
        ('- 날짜 포맷: YYYY.MM.DD (HH:mm)', CODE_FG),
        ('', None),
        ('... 계속', CM),
    ]
    ly = cy + 24
    for text, color in code_lines:
        if color is None:
            ly += 18
            continue
        add_text(s, cx + 28, ly, cw - 56, 22, text, size=12, color=color,
                 font=FONT_MONO)
        ly += 22
    # 353 lines chip
    add_rect(s, cx + cw - 140, cy + 20, 120, 34, fill=INK, radius=6)
    add_text(s, cx + cw - 140, cy + 26, 120, 24, '353 LINES', size=12, bold=True,
             color=YELLOW, letter_spacing=0.14, align='center')
    return s

def s19(prs):
    s = new_slide(prs, PAPER)
    add_top_bar(s, "§ 3 · STEP 2", "EMILY · SELFISH CLUB · AAA")
    add_kicker(s, 140, 180, "그래서 그냥 걔한테 물어봤어요.")
    add_text(s, 140, 260, 1640, 40, "TO: CLAUDE", size=14, bold=True,
             color=RGBColor(0x99, 0x99, 0x99), letter_spacing=0.2)
    add_runs(s, 140, 330, 1640, 500, [
        {'text': '"너한테 기억 잘 주려면\n뭘 가져야 해?\n…네가 생각하는 ',
         'size': 88, 'bold': True, 'color': INK, 'letter_spacing': -0.04},
        {'text': '하네스', 'size': 88, 'bold': True, 'color': INK, 'highlight': YELLOW,
         'letter_spacing': -0.04},
        {'text': '는 뭐야?"', 'size': 88, 'bold': True, 'color': INK,
         'letter_spacing': -0.04},
    ], line_height=1.15)
    add_text(s, 140, 900, 1640, 40,
             "AAA팀에서 '하네스'라는 단어만 언뜻 들었던 참이었거든요.",
             size=20, color=RGBColor(0x66, 0x66, 0x66))
    return s

def s20(prs):
    s = new_slide(prs, PAPER)
    add_top_bar(s, "§ 3 · STEP 2", "EMILY · SELFISH CLUB · AAA")
    add_kicker(s, 140, 160, "걔가 답해준 거 정리해보니")
    add_text(s, 140, 210, 1640, 100, "세 가지로 모이더라고요.", size=48, bold=True,
             color=INK, letter_spacing=-0.035)
    cards = [
        ('01', 'SPACE',    '공간', '기록이 쌓일 자리를 만들어주기.', '폴더 구조 · 파일 네이밍 규칙'),
        ('02', 'MEMORY',   '기억', '세션 시작할 때 자동으로 읽게 만들기.', 'MEMORY.md · 자동 로드'),
        ('03', 'LEARNING', '학습', '규칙을 안 지키면 막게 만들기.', 'pre-commit 훅 · 커밋 차단'),
    ]
    x0 = 140
    cw = 540
    gap = 20
    y0 = 380
    for i, (n, k, ko, t, d) in enumerate(cards):
        x = x0 + i * (cw + gap)
        add_rect(s, x, y0, cw, 480, fill=SURFACE, radius=16)
        add_text(s, x + 40, y0 + 40, cw - 80, 30, f'{n} / {k}', size=14, bold=True,
                 color=RGBColor(0xAA, 0xAA, 0xAA))
        add_text(s, x + 40, y0 + 90, cw - 80, 80, ko, size=44, bold=True, color=INK)
        add_text(s, x + 40, y0 + 200, cw - 80, 140, t, size=22, bold=True, color=INK,
                 line_height=1.4)
        add_text(s, x + 40, y0 + 380, cw - 80, 60, d, size=14,
                 color=RGBColor(0x77, 0x77, 0x77))
    add_runs(s, 140, 900, 1640, 40, [
        {'text': '적용해보니 잘 돌더라고요. 알고 보니 이게 ', 'size': 20, 'color': INK},
        {'text': "'하네스'", 'size': 20, 'bold': True, 'color': INK},
        {'text': '였던 거예요.', 'size': 20, 'color': INK},
    ])
    return s

def s21(prs):
    s = new_slide(prs, INK)
    add_top_bar(s, "§ 3 · STEP 2", "EMILY · SELFISH CLUB · AAA", dark=True)
    add_text(s, 140, 180, 1640, 40, "근데 이게 말로만 들으면 감이 안 와요.", size=18,
             bold=True, color=RGBColor(0x99, 0x99, 0x99), letter_spacing=0.22)
    add_runs(s, 140, 280, 1640, 500, [
        {'text': '실제로는 이렇게\n', 'size': 130, 'bold': True, 'color': PAPER,
         'letter_spacing': -0.04},
        {'text': '작동', 'size': 130, 'bold': True, 'color': INK, 'highlight': YELLOW,
         'letter_spacing': -0.04},
        {'text': '해요.', 'size': 130, 'bold': True, 'color': PAPER,
         'letter_spacing': -0.04},
    ], line_height=1.02)
    add_text(s, 140, 900, 1640, 40,
             'MEMORY.md 인덱스 갱신 — 한 장면에 3축이 다 걸려 있어요.', size=22,
             color=RGBColor(0xBB, 0xBB, 0xBB))
    return s

def s22(prs):
    s = new_slide(prs, PAPER)
    add_top_bar(s, "§ 3 · STEP 2", "EMILY · SELFISH CLUB · AAA")
    add_kicker(s, 140, 140, "세션 저장 → 커밋 차단 → 인덱스 추가.")
    add_runs(s, 140, 190, 1640, 120, [
        {'text': '이 한 ', 'size': 40, 'bold': True, 'color': INK,
         'letter_spacing': -0.035},
        {'text': '사이클', 'size': 40, 'bold': True, 'color': INK, 'highlight': YELLOW,
         'letter_spacing': -0.035},
        {'text': '이 곧 하네스가 돌아가는 모습이에요.', 'size': 40, 'bold': True,
         'color': INK, 'letter_spacing': -0.035},
    ])
    # 왼쪽 박스 (폴더 구조)
    lx, ly, lw, lh = 140, 340, 830, 480
    add_rect(s, lx, ly, lw, lh, fill=RGBColor(0xF8, 0xF8, 0xF8), radius=16)
    add_text(s, lx + 28, ly + 28, lw - 56, 30, '공간 · 폴더 구조', size=12, bold=True,
             color=RGBColor(0x88, 0x88, 0x88), letter_spacing=0.2)
    folder_lines = [
        '📁 memory/',
        '├─ MEMORY.md   ← 인덱스',
        '├─ session_0418_aaa.md',
        '├─ session_0421_crm.md  [new]',
        '├─ feedback_0419.md',
        '└─ ...',
        '',
        '📁 .husky/',
        '└─ pre-commit   ← 검사',
    ]
    cy = ly + 70
    for line in folder_lines:
        add_text(s, lx + 28, cy, lw - 56, 24, line, size=12, color=RGBColor(0x22, 0x22, 0x22),
                 font=FONT_MONO)
        cy += 28
    add_text(s, lx + 28, ly + lh - 70, lw - 56, 40,
             '세션 파일이 쌓이면 MEMORY.md에도 한 줄 추가돼야 함.',
             size=12, color=RGBColor(0x66, 0x66, 0x66))
    # 오른쪽 박스 (터미널)
    rx, ry, rw, rh = 990, 340, 790, 480
    add_rect(s, rx, ry, rw, rh, fill=CODE_BG, radius=16)
    add_text(s, rx + 28, ry + 28, rw - 56, 30, '학습 · PRE-COMMIT 훅이 막음', size=12,
             bold=True, color=YELLOW, letter_spacing=0.2)
    term_lines = [
        ('$', RGBColor(0xAA, 0xAA, 0xAA), ' git commit -m "add session_0421"', CODE_FG),
        ('', None, '', None),
        ('running pre-commit…', RGBColor(0x88, 0x88, 0x88)),
        ('✓ lint', RGBColor(0x7E, 0xCC, 0x7A)),
        ('✓ type-check', RGBColor(0x7E, 0xCC, 0x7A)),
        ('✗ memory-index', RGBColor(0xFF, 0x5A, 0x5A)),
        ('', None),
        ('❌ 인덱스 누락 1건 — MEMORY.md에 추가 필요', RGBColor(0xFF, 0x5A, 0x5A)),
        ('· memory/session_0421_crm.md', RGBColor(0xBB, 0xBB, 0xBB)),
        ('', None),
        ('commit aborted.', RGBColor(0xFF, 0x5A, 0x5A)),
    ]
    cy = ry + 80
    for parts in term_lines:
        if len(parts) == 2 and parts[1] is None:
            cy += 12
            continue
        runs = []
        for i in range(0, len(parts), 2):
            if parts[i + 1] is None: continue
            runs.append({'text': parts[i], 'size': 12, 'color': parts[i + 1],
                         'font': FONT_MONO})
        if runs:
            add_runs(s, rx + 28, cy, rw - 56, 26, runs)
        cy += 24
    # 하단 옐로우 바 3개
    tags = [
        ('SPACE',    '폴더 구조가 자리를 만듦'),
        ('MEMORY',   'MEMORY.md가 매 세션 자동 로드'),
        ('LEARNING', '훅이 규칙을 강제함'),
    ]
    tx = 140
    tw = 540
    tgap = 20
    ty = 870
    for i, (k, d) in enumerate(tags):
        x = tx + i * (tw + tgap)
        add_rect(s, x, ty, tw, 80, fill=YELLOW, radius=10)
        add_text(s, x + 20, ty + 12, tw - 40, 24, k, size=12, bold=True, color=YELLOW_INK,
                 letter_spacing=0.18)
        add_text(s, x + 20, ty + 40, tw - 40, 40, d, size=16, bold=True, color=INK)
    return s

def s23(prs):
    s = new_slide(prs, PAPER)
    add_top_bar(s, "§ 4 · 유즈케이스", "EMILY · SELFISH CLUB · AAA")
    add_rect(s, 140, 160, 170, 34, fill=INK, radius=4)
    add_text(s, 140, 166, 170, 24, 'USE CASE ①', size=12, bold=True, color=YELLOW,
             letter_spacing=0.18, align='center')
    add_runs(s, 140, 220, 1640, 140, [
        {'text': '규칙을 ', 'size': 52, 'bold': True, 'color': INK,
         'letter_spacing': -0.035},
        {'text': 'md 파일', 'size': 52, 'bold': True, 'color': INK, 'highlight': YELLOW,
         'letter_spacing': -0.035},
        {'text': '로 빼서 AI가 읽게.', 'size': 52, 'bold': True, 'color': INK,
         'letter_spacing': -0.035},
    ])
    add_text(s, 140, 370, 1640, 50, '353줄 CLAUDE.md → 역할별 PRD로 분리.', size=22,
             color=RGBColor(0x66, 0x66, 0x66))
    # BEFORE
    bx, by, bw, bh = 140, 460, 760, 420
    add_rect(s, bx, by, bw, bh, fill=SURFACE, radius=16)
    add_rect(s, bx + 32, by + 30, 130, 30, fill=RGBColor(0xDD, 0xDD, 0xDD), radius=4)
    add_text(s, bx + 32, by + 35, 130, 24, 'BEFORE', size=11, bold=True,
             color=RGBColor(0x44, 0x44, 0x44), letter_spacing=0.18, align='center')
    add_text(s, bx + 32, by + 80, bw - 64, 50, 'CLAUDE.md 하나에 다 박음', size=22,
             bold=True, color=INK)
    add_rect(s, bx + 32, by + 150, bw - 64, 190, fill=PAPER, line=RULE, line_w=1, radius=10)
    add_text(s, bx + 52, by + 170, bw - 104, 40, '📄 CLAUDE.md  353 lines', size=16,
             bold=True, color=INK)
    add_text(s, bx + 52, by + 210, bw - 104, 120,
             '톤 · 타이밍 · 타겟 · 링크 ·\n변수 · 예외 · 채널별 규칙…\n(전부 한 파일)',
             size=13, color=RGBColor(0x55, 0x55, 0x55), font=FONT_MONO, line_height=1.7)
    add_text(s, bx + 32, by + 360, bw - 64, 40, '→ 길어서 AI가 끝까지 안 읽음', size=14,
             color=RGBColor(0xC4, 0x34, 0x3B))
    # ARROW
    add_text(s, 910, 640, 100, 80, '→', size=56, color=RGBColor(0xBB, 0xBB, 0xBB),
             align='center')
    # AFTER
    ax, ay, aw, ah = 1020, 460, 760, 420
    add_rect(s, ax, ay, aw, ah, fill=INK, radius=16)
    add_rect(s, ax + 32, ay + 30, 130, 30, fill=YELLOW, radius=4)
    add_text(s, ax + 32, ay + 35, 130, 24, 'AFTER', size=11, bold=True, color=INK,
             letter_spacing=0.18, align='center')
    add_text(s, ax + 32, ay + 80, aw - 64, 50, '역할별 PRD로 쪼갬', size=22, bold=True,
             color=PAPER)
    add_rect(s, ax + 32, ay + 150, aw - 64, 190, fill=RGBColor(0x1E, 0x1E, 0x1E),
             radius=10)
    prd_lines = [
        '📁 prd/',
        '├─ alimtalk.md   알림톡 규칙',
        '├─ tone.md         톤앤매너',
        '├─ targeting.md  예외 규칙',
        '└─ links.md        UTM · 단축',
    ]
    py = ay + 170
    for line in prd_lines:
        add_text(s, ax + 52, py, aw - 104, 26, line, size=13, color=YELLOW,
                 font=FONT_MONO)
        py += 28
    add_text(s, ax + 32, ay + 360, aw - 64, 40, '→ AI가 그때그때 필요한 것만 로드',
             size=14, color=YELLOW)
    return s

def s24(prs):
    s = new_slide(prs, PAPER)
    add_top_bar(s, "§ 4 · 유즈케이스", "EMILY · SELFISH CLUB · AAA")
    add_rect(s, 140, 160, 170, 34, fill=INK, radius=4)
    add_text(s, 140, 166, 170, 24, 'USE CASE ②', size=12, bold=True, color=YELLOW,
             letter_spacing=0.18, align='center')
    add_runs(s, 140, 220, 1640, 140, [
        {'text': '에이전트 8개', 'size': 52, 'bold': True, 'color': INK,
         'highlight': YELLOW, 'letter_spacing': -0.035},
        {'text': ', 역할을 쪼갰어요.', 'size': 52, 'bold': True, 'color': INK,
         'letter_spacing': -0.035},
    ])
    add_text(s, 140, 370, 1640, 50, '한 명이 다 하는 대신, 직원 여덟 명처럼.', size=20,
             color=RGBColor(0x66, 0x66, 0x66))
    agents = [
        ('01', 'director',        '전체 흐름 지휘', False),
        ('02', 'copywriter',      '카피 초안 작성', False),
        ('03', 'data-collector',  'DB에서 타겟 조회', False),
        ('04', 'dispatcher',      '실제 발송 트리거', False),
        ('05', 'verifier',        '링크·변수 검수', False),
        ('06', 'scheduler',       '예약 시점 관리', False),
        ('07', 'reporter',        '발송 리포트', False),
        ('08', 'code-reviewer',   '검수 게이트', True),
    ]
    x0 = 140
    cw = 400
    gap = 14
    y0 = 440
    rh = 180
    for i, (n, t, d, hl) in enumerate(agents):
        r, c = i // 4, i % 4
        x = x0 + c * (cw + gap)
        y = y0 + r * (rh + 14)
        fill = YELLOW if hl else SURFACE
        fg = YELLOW_INK if hl else RGBColor(0x99, 0x99, 0x99)
        tc = INK if not hl else INK
        dc = YELLOW_INK if hl else RGBColor(0x66, 0x66, 0x66)
        add_rect(s, x, y, cw, rh, fill=fill, radius=16)
        add_text(s, x + 28, y + 24, cw - 56, 30, n, size=14, bold=True, color=fg)
        add_text(s, x + 28, y + 60, cw - 56, 50, t, size=20, bold=True, color=tc,
                 font=FONT_MONO)
        add_text(s, x + 28, y + 120, cw - 56, 40, d, size=14, color=dc)
    add_text(s, 140, 850, 1640, 60,
             '↑ code-reviewer는 AAA 과제 중에 \'검수가 필요하다\'는 걸 깨닫고 뒤늦게 추가한 역할이에요.',
             size=16, color=RGBColor(0x66, 0x66, 0x66))
    return s

def s25(prs):
    s = new_slide(prs, PAPER)
    add_top_bar(s, "§ 4 · 유즈케이스", "EMILY · SELFISH CLUB · AAA")
    add_rect(s, 140, 160, 170, 34, fill=INK, radius=4)
    add_text(s, 140, 166, 170, 24, 'USE CASE ③', size=12, bold=True, color=YELLOW,
             letter_spacing=0.18, align='center')
    add_runs(s, 140, 220, 1640, 140, [
        {'text': '발송이 ', 'size': 48, 'bold': True, 'color': INK,
         'letter_spacing': -0.035},
        {'text': '다음 발송', 'size': 48, 'bold': True, 'color': INK, 'highlight': YELLOW,
         'letter_spacing': -0.035},
        {'text': '을 만들어요.', 'size': 48, 'bold': True, 'color': INK,
         'letter_spacing': -0.035},
    ])
    add_text(s, 140, 370, 1640, 50, '알림톡 한 번 나가면, 그다음이 자동으로 세팅돼요.',
             size=20, color=RGBColor(0x66, 0x66, 0x66))
    # 가로 플로우
    stages = [
        (280, 'STEP 1 · 내가 누름',  '알림톡 ①\n발송',          YELLOW, INK, YELLOW_INK),
        (380, 'AGENT',              'data-collector\n신청자 조회', RGBColor(0xF3, 0xF3, 0xF3), INK, RGBColor(0x88, 0x88, 0x88)),
        (380, 'AGENT',              'dispatcher\n신청자 제외 세팅', RGBColor(0xF3, 0xF3, 0xF3), INK, RGBColor(0x88, 0x88, 0x88)),
        (300, 'STEP 2 · 자동',      '미신청자 전용\n알림톡 ② 발송', INK, PAPER, YELLOW),
    ]
    fx = 140
    fy = 450
    fh = 200
    for i, (w, tag, t, bg, fg, tag_fg) in enumerate(stages):
        add_rect(s, fx, fy, w, fh, fill=bg, radius=12)
        add_text(s, fx + 22, fy + 22, w - 44, 24, tag, size=10, bold=True, color=tag_fg,
                 letter_spacing=0.16)
        add_text(s, fx + 22, fy + 60, w - 44, 120, t, size=18, bold=True, color=fg,
                 line_height=1.3)
        fx += w + 10
        if i < len(stages) - 1:
            add_text(s, fx, fy + 70, 40, 60, '→', size=32, color=RGBColor(0xBB, 0xBB, 0xBB),
                     align='center')
            fx += 40 + 10
    # 아래 요약 박스
    add_rect(s, 140, 680, 1640, 80, fill=INK, radius=10)
    add_runs(s, 170, 710, 1580, 40, [
        {'text': 'STEP 1만 내가 승인하면 돼요', 'size': 18, 'bold': True, 'color': YELLOW},
        {'text': ' — 나머지는 에이전트끼리 릴레이로 STEP 2를 만들어요.', 'size': 18,
         'color': PAPER},
    ])
    # 옛날엔 / 지금은
    add_text(s, 140, 790, 820, 100,
             '옛날엔 알림톡 ① 보내고 → 다음 날 수동으로 신청자 빼고 → 미신청자만 추려서 ② 보냄.',
             size=14, color=RGBColor(0x66, 0x66, 0x66), line_height=1.5)
    add_text(s, 960, 790, 820, 100,
             '지금은 알림톡 ① 발송 직후에 에이전트끼리 자동으로 다음 타겟을 재조회.',
             size=14, color=RGBColor(0x66, 0x66, 0x66), line_height=1.5)
    return s

def s26(prs):
    s = new_slide(prs, PAPER)
    add_top_bar(s, "§ 5 · IMPACT", "EMILY · SELFISH CLUB · AAA")
    add_kicker(s, 140, 160, "그래서 뭐가 바뀌었냐면.")
    rows = [
        ('알림톡 1건', '3h', '5min', '매일 아침 슬랙 모달 한 번.', False),
        ('공유회 준비', '2d', '1h',   '채널별 발송 세팅 전체.',   False),
        ('리스크',      '오타·링크 실수', '검수 단계에서 차단.',       '', True),
    ]
    top = 260
    rh = 180
    for i, (lab, a, b, note, long) in enumerate(rows):
        y = top + i * rh
        add_rect(s, 140, y, 1640, 2, fill=RULE)
        add_text(s, 140, y + 60, 220, 40, lab.upper(), size=14, bold=True,
                 color=RGBColor(0x88, 0x88, 0x88), letter_spacing=0.16)
        asize = 32 if long else 80
        add_text(s, 380, y + 20, 400, 140, a, size=asize, bold=True,
                 color=RGBColor(0x99, 0x99, 0x99), letter_spacing=-0.03)
        add_text(s, 820, y + 40, 80, 80, '→', size=40,
                 color=RGBColor(0xBB, 0xBB, 0xBB), align='center')
        bsize = 32 if long else 80
        # hl b
        add_runs(s, 940, y + 20, 500, 140, [
            {'text': b, 'size': bsize, 'bold': True, 'color': INK, 'highlight': YELLOW,
             'letter_spacing': -0.03},
        ])
        add_text(s, 1460, y + 60, 320, 80, note, size=14,
                 color=RGBColor(0x55, 0x55, 0x55))
    # 하단 메시지
    add_runs(s, 140, 870, 1640, 80, [
        {'text': '퇴근 직전 카피 쓰다 지쳤던 제가,\n', 'size': 26, 'bold': True,
         'color': INK},
        {'text': '맥북 닫고 집에 가도 메시지가 나가요.', 'size': 26, 'bold': True,
         'color': INK, 'highlight': YELLOW},
    ], line_height=1.4)
    return s

def s27(prs):
    s = new_slide(prs, PAPER)
    add_top_bar(s, "§ 6 · 팀", "EMILY · SELFISH CLUB · AAA")
    add_kicker(s, 140, 160, "혼자 한 거 아니에요.")
    add_text(s, 140, 210, 1640, 100, 'AAA 6주, 옆자리에서 이렇게 영향받았어요.',
             size=44, bold=True, color=INK, letter_spacing=-0.035)
    people = [
        ('비비안', '전체 대시보드 크루랑 협업하면서,\n내 CRM이 전체 어디에 들어가는지 감 잡음.',
         False),
        ('다니', '마케팅 OS 보면서 영향받아,\nCLAUDE.md → PRD 구조 차용.', False),
        ('code-reviewer',
         "AAA 과제로 하네스 서로 공유하다\n'검수가 필요하다' 깨닫고 만든 에이전트.", True),
        ('흐민 · 오웬', "흐민의 '질문하는 AI', 오웬의 '찜마켓'\n보면서 활용 범위 확장.", False),
    ]
    x0 = 140
    cw = 400
    gap = 14
    y0 = 400
    for i, (who, t, mono) in enumerate(people):
        x = x0 + i * (cw + gap)
        add_rect(s, x, y0, cw, 280, fill=SURFACE, radius=16)
        add_text(s, x + 28, y0 + 28, cw - 56, 40, who, size=18, bold=True, color=INK,
                 font=FONT_MONO if mono else FONT_KO)
        add_rect(s, x + 28, y0 + 84, cw - 56, 1, fill=RULE)
        add_text(s, x + 28, y0 + 110, cw - 56, 160, t, size=14,
                 color=RGBColor(0x33, 0x33, 0x33), line_height=1.55)
    # 하단 옐로우 박스
    add_rect(s, 140, 760, 1640, 160, fill=YELLOW, radius=14)
    add_runs(s, 180, 790, 1560, 60, [
        {'text': '그래서 ', 'size': 32, 'bold': True, 'color': INK},
        {'text': '스폰지클럽', 'size': 32, 'bold': True, 'color': INK, 'underline': True},
        {'text': '.', 'size': 32, 'bold': True, 'color': INK},
    ])
    add_text(s, 180, 850, 1560, 60, '빨아들이고, 내보내고, 또 누군가가 빨아들이고.',
             size=18, color=YELLOW_INK)
    return s

def s28(prs):
    s = new_slide(prs, INK)
    add_top_bar(s, "OUTRO", "EMILY · SELFISH CLUB · AAA", dark=True)
    add_bot_bar(s, "에밀리 · 셀피쉬클럽 CRM PM", "AAA 공유회 · 2026.04.28", dark=True)
    add_text(s, 140, 180, 1640, 40, "여러분은 이렇게 해보세요.", size=18, bold=True,
             color=RGBColor(0x99, 0x99, 0x99), letter_spacing=0.22)
    add_runs(s, 140, 240, 1640, 140, [
        {'text': '내일부터 ', 'size': 78, 'bold': True, 'color': PAPER,
         'letter_spacing': -0.04},
        {'text': '세 가지', 'size': 78, 'bold': True, 'color': INK, 'highlight': YELLOW,
         'letter_spacing': -0.04},
        {'text': '만.', 'size': 78, 'bold': True, 'color': PAPER,
         'letter_spacing': -0.04},
    ])
    rows = [
        ('01', '한 파일로 시작하기.',      'CLAUDE.md에 내 규칙부터 적어보기.'),
        ('02', '걔한테 먼저 물어보기.',    '"내 맥락 잘 기억하려면 뭘 줘야 해?"'),
        ('03', '사람 관문 하나 남겨두기.', '완전 자동 말고, 승인 한 번.'),
    ]
    ry = 440
    for n, t, d in rows:
        add_text(s, 140, ry, 100, 60, n, size=32, bold=True, color=YELLOW)
        add_text(s, 260, ry, 1500, 50, t, size=30, bold=True, color=PAPER)
        add_text(s, 260, ry + 44, 1500, 40, d, size=18,
                 color=RGBColor(0xBB, 0xBB, 0xBB))
        ry += 100
    add_rect(s, 140, 820, 1640, 2, fill=RGBColor(0x33, 0x33, 0x33))
    add_runs(s, 140, 860, 1640, 120, [
        {'text': '감사합니다. ', 'size': 80, 'bold': True, 'color': PAPER,
         'letter_spacing': -0.04},
        {'text': 'Q & A.', 'size': 80, 'bold': True, 'color': INK, 'highlight': YELLOW,
         'letter_spacing': -0.04},
    ])
    return s

# ── Main ───────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W * 9525)
    prs.slide_height = Emu(SLIDE_H * 9525)
    builders = [s01, s02, s03, s04, s05, s06, s07, s08, s09, s10,
                s11, s12, s13, s14, s15, s16, s17, s18, s19, s20,
                s21, s22, s23, s24, s25, s26, s27, s28]
    for i, b in enumerate(builders, 1):
        print(f"[{i:02d}/28] {b.__name__} ...")
        b(prs)
    prs.save(OUT)
    print(f"\n✓ saved: {OUT.name}")
    print(f"  size: {OUT.stat().st_size / 1024:.0f}KB")

if __name__ == "__main__":
    main()
