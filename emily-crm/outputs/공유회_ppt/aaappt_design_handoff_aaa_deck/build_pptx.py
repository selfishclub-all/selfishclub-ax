"""
AAA 발표 덱 v7/v8 → PPTX 변환기
전략: 엔트리 HTML + CSS/JS/JSX를 인라인해 standalone 생성 → Chrome 헤드리스 스크린샷 N번 → pptx 조립
"""
import subprocess, sys
from pathlib import Path

HERE = Path(__file__).parent
ENTRY = HERE / "AAA 발표 덱 v7.html"
STANDALONE = HERE / "_standalone_v8.html"
PATCHED = HERE / "_standalone_patched.html"
OUT_PNG_DIR = HERE / "_slide_pngs"
OUT_PPTX = HERE / "AAA_발표_덱_v7.pptx"
CHROME = "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"

SLIDE_W = 1920
SLIDE_H = 1080
CAPTURE_H = 1200
TOTAL_SLIDES = 28  # 1단계: 기존 번호 유지. 2단계에서 30으로 확장 예정.

JSX_FILES = ["slides-core", "slides-1", "slides-2", "slides-3", "slides-4", "slides-5"]

INJECT = """
<script>
(function(){
  function go(){
    var m = location.search.match(/slide=(\\d+)/);
    if (!m) return;
    var idx = parseInt(m[1], 10);
    var ds = document.querySelector('deck-stage');
    if (!ds || typeof ds.goTo !== 'function'){ setTimeout(go, 100); return; }
    ds.goTo(idx);
    setTimeout(function(){ document.body.setAttribute('data-ready', '1'); }, 800);
  }
  window.addEventListener('load', function(){ setTimeout(go, 500); });
})();
</script>
"""

def build_standalone():
    """엔트리 HTML + CSS/JS/JSX 인라인 → standalone 파일로 출력 (assets/는 상대경로 유지)"""
    html = ENTRY.read_text(encoding="utf-8")
    # CSS 인라인
    css = (HERE / "deck-styles.css").read_text(encoding="utf-8")
    html = html.replace('<link rel="stylesheet" href="deck-styles.css">',
                        f'<style>\n{css}\n</style>')
    # deck-stage.js 인라인
    ds = (HERE / "deck-stage.js").read_text(encoding="utf-8")
    html = html.replace('<script src="deck-stage.js"></script>',
                        f'<script>\n{ds}\n</script>')
    # JSX 인라인
    for name in JSX_FILES:
        jsx = (HERE / f"{name}.jsx").read_text(encoding="utf-8")
        # </script> 이스케이프
        jsx = jsx.replace('</script>', '<\\/script>')
        old = f'<script type="text/babel" src="{name}.jsx"></script>'
        new = f'<script type="text/babel">\n{jsx}\n</script>'
        html = html.replace(old, new)
    STANDALONE.write_text(html, encoding="utf-8")
    print(f"✓ built standalone → {STANDALONE.name} ({STANDALONE.stat().st_size//1024}KB)")

def patch_html():
    html = STANDALONE.read_text(encoding="utf-8")
    patched = html.replace("</body>", INJECT + "\n</body>")
    PATCHED.write_text(patched, encoding="utf-8")
    print(f"✓ patched → {PATCHED.name} ({PATCHED.stat().st_size//1024}KB)")

def postprocess_png(png_path):
    from PIL import Image
    img = Image.open(png_path).convert("RGB")
    w, h = img.size
    px = img.load()
    def is_black_row(y, thr=30):
        return all(sum(px[x, y]) < thr for x in range(0, w, 40))
    top = 0
    while top < h and is_black_row(top):
        top += 1
    bot = min(h, top + 1080)
    cropped = img.crop((0, top, w, bot))
    if cropped.size != (1920, 1080):
        cropped = cropped.resize((1920, 1080), Image.LANCZOS)
    cropped.save(png_path)

def capture_slides():
    OUT_PNG_DIR.mkdir(exist_ok=True)
    for i in range(TOTAL_SLIDES):
        png = OUT_PNG_DIR / f"slide-{i+1:02d}.png"
        url = f"file://{PATCHED}?slide={i}"
        cmd = [
            CHROME, "--headless=new", "--disable-gpu", "--no-sandbox", "--hide-scrollbars",
            f"--window-size={SLIDE_W},{CAPTURE_H}",
            f"--screenshot={png}",
            "--virtual-time-budget=15000",
            "--timeout=30000",
            url,
        ]
        print(f"[{i+1:02d}/{TOTAL_SLIDES}] capturing...", flush=True)
        r = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
        if not png.exists() or png.stat().st_size < 5000:
            print(f"  ⚠️ fail ({r.returncode}): {r.stderr[:300]}")
        else:
            postprocess_png(png)
            print(f"  ✓ {png.stat().st_size//1024}KB")

def build_pptx():
    from pptx import Presentation
    from pptx.util import Emu
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Emu(1920 * 9525)
    prs.slide_height = Emu(1080 * 9525)
    blank = prs.slide_layouts[6]
    count = 0
    for i in range(1, TOTAL_SLIDES + 1):
        png = OUT_PNG_DIR / f"slide-{i:02d}.png"
        if not png.exists():
            print(f"  skip {i}: no png")
            continue
        slide = prs.slides.add_slide(blank)
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        slide.shapes.add_picture(
            str(png),
            Emu(-1 * 9525), Emu(-1 * 9525),
            width=Emu((1920 + 2) * 9525),
            height=Emu((1080 + 2) * 9525)
        )
        count += 1
    prs.save(OUT_PPTX)
    print(f"\n✓ saved: {OUT_PPTX.name}")
    print(f"  slides: {count}  size: {OUT_PPTX.stat().st_size/1024/1024:.1f}MB")

if __name__ == "__main__":
    mode = sys.argv[1] if len(sys.argv) > 1 else "all"
    if mode in ("all", "build"):   build_standalone()
    if mode in ("all", "patch"):   patch_html()
    if mode in ("all", "capture"): capture_slides()
    if mode in ("all", "pptx"):    build_pptx()
