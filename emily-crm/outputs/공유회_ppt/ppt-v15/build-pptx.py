#!/usr/bin/env python3
"""Assemble 31 PNG screenshots into a single 16:9 PPTX deck."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

HERE = Path(__file__).parent
PNG_DIR = HERE / "png-export"
OUT = HERE / "AAA_공유회_v15.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9 widescreen (1920px @ 144dpi)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]

pngs = sorted(PNG_DIR.glob("slide-*.png"))
assert len(pngs) == 31, f"Expected 31 PNGs, got {len(pngs)}"

for png in pngs:
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.add_picture(
        str(png), 0, 0,
        width=prs.slide_width, height=prs.slide_height
    )

prs.save(OUT)
print(f"[done] {len(pngs)} slides → {OUT}")
print(f"[size] {OUT.stat().st_size / 1024 / 1024:.1f} MB")
